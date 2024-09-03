import argparse
import os
import re
import shutil
import traceback
import zipfile

import cv2
import imagehash
import markdown
import markdown_it
import torch
from moviepy.editor import VideoFileClip
from PIL import Image
from transformers import (AutoModelForCausalLM, AutoModelForSpeechSeq2Seq,
                          AutoProcessor, AutoTokenizer, pipeline)
# from optimum.nvidia import AutoModelForCausalLM
# from optimum.nvidia.pipelines import pipeline
# from transformers import AutoProcessor, AutoTokenizer, AutoModelForSpeechSeq2Seq

DEFAULT_SYSTEM_PROMPT = "You are a helpful assistant."
DEFAULT_SUBTITLE_SUMMARY_PROMPT = "下面是视频提取的部分字幕片段，该片段对应的是一页PPT内容，该视频是学习类的视频，请你帮我总结改页PPT对应视频字幕的主要内容，使用markdown" \
                                  "格式进行输出，请用中文回答"
DEFAULT_VIDEO_SUMMARY_PROMPT = "以下是视频每一页ppt的总结，请你帮我总结整个视频的主要内容，生成总结性的内容，使用markdown格式进行输出，请用中文回答"
DEFAULT_VIDEO_MINDMAP_PROMPT = "以下是视频每一页ppt的总结，请你帮我尝试总结整个视频的主要内容，生成思维导图，请使用markdown格式进行输出该思维导图，请用中文回答"

device = "cuda" if torch.cuda.is_available() else "cpu"


def extract_similar_frames(video_path, output_extract_folder, interval_sec=1.5, threshold=5, similarity_algo="ahash"):
    output_frame_path = os.path.join(output_extract_folder, "frames")
    if not os.path.exists(output_frame_path):
        os.makedirs(output_frame_path)
    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        raise ValueError(f"无法打开视频文件: {video_path}")
    fps = cap.get(cv2.CAP_PROP_FPS)
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    interval = int(fps * interval_sec)

    frame_count = 0
    image_extract_dict = {}
    try:
        while True:
            ret, frame = cap.read()
            if not ret:
                break

            if frame_count % interval == 0:
                image = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
                extract_first_frame(image, frame_count, image_extract_dict, output_frame_path,
                                    threshold, similarity_algo)  # Assuming defined elsewhere

            # Progress update
            if frame_count % (total_frames // 100) == 0:  # Update progress every 1%
                progress = (frame_count / total_frames) * 100
                print(f"Processing progress: {progress:.2f}%")
            frame_count += 1
    finally:
        cap.release()
        print("Frame extraction completed.")
    return list_similar_frames(output_frame_path)


def list_similar_frames(frames_dir):
    result = []
    image_files = sorted([f for f in os.listdir(frames_dir) if f.endswith((".jpg", ".png", ".jpeg"))],
                         key=lambda x: int(os.path.splitext(x)[0]))
    for i in image_files:
        result.append((os.path.join(frames_dir, i), f"frame {i}"))
    return result


def extract_first_frame(image, frame_num, image_extract_dict, output_frame_path, threshold=5, similarity_algo="ahash"):
    if similarity_algo == "ahash":
        image_hash = imagehash.average_hash(image)
    elif similarity_algo == "phash":
        image_hash = imagehash.phash(image)
    elif similarity_algo == "dhash":
        image_hash = imagehash.dhash(image)
    elif similarity_algo == "whash":
        image_hash = imagehash.whash(image)
    else:
        raise ValueError("Invalid similarity algorithm")
    for _, image_info in image_extract_dict.items():
        if image_hash - image_info["hash"] <= threshold:
            return  # 发现相似图片，不再保存

    frame_filename = os.path.join(output_frame_path, f"{frame_num}.jpg")
    image_extract_dict[frame_filename] = {
        "hash": image_hash,
        "frame": frame_num
    }
    image.save(frame_filename)


def format_time(seconds):
    """ Helper function to convert seconds to time format (H:M:S) """
    hours = int(seconds // 3600)
    minutes = int((seconds % 3600) // 60)
    seconds = int(seconds % 60)  # Convert to integer to avoid decimals
    return f"{hours:02}:{minutes:02}:{seconds:02}"


def calculate_time_segments(video_file, output_extract_folder):
    output_frame_path = os.path.join(output_extract_folder, "frames")
    frame_filenames = [f for f in os.listdir(output_frame_path) if f.endswith((".jpg", ".png", ".jpeg"))]
    cap = cv2.VideoCapture(video_file)
    if not cap.isOpened():
        raise ValueError("Cannot open video file")

    fps = cap.get(cv2.CAP_PROP_FPS)
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    video_duration = total_frames // fps
    cap.release()

    frame_numbers = sorted(int(fname.split('.')[0]) for fname in frame_filenames)
    time_segments = []
    for i, frame_num in enumerate(frame_numbers):
        start_time = frame_num // fps
        if i + 1 < len(frame_numbers):
            end_time = (frame_numbers[i + 1] - 1) // fps
        else:
            end_time = video_duration  # Use total video duration for the last segment
        print(f"split video {format_time(start_time)}--{format_time(end_time)} segments")
        time_segments.append((start_time, end_time))
    return time_segments


def extract_relative_subtitle_with_merge(video_path, output_extract_folder, merge=True, file_size_threshold=1024):
    extract_relative_subtitle(video_path, output_extract_folder)
    if merge:
        merge_small_subtitles(video_path, output_extract_folder, file_size_threshold)
    return list_all_subtitles(os.path.join(output_extract_folder))


def list_all_subtitles(output_extract_folder):
    subtitles_dir = os.path.join(output_extract_folder, "subtitles")
    print("list_all_subtitles", subtitles_dir)
    files = []
    subtitle_files = sorted([f for f in os.listdir(subtitles_dir) if f.endswith(".txt")],
                            key=lambda x: int(os.path.splitext(x)[0]))

    for i in subtitle_files:
        files.append(os.path.join(subtitles_dir, i))
    return files


def list_all_summary_files(output_extract_folder):
    summary_dir = os.path.join(output_extract_folder, "summary")
    print("list_all_summary_files", summary_dir)
    files = []
    summary_files = sorted(
        [f for f in os.listdir(summary_dir) if f.endswith(".md") and f != "summary.md" and f != "mindmap.md"],
        key=lambda x: int(os.path.splitext(x)[0]))

    for i in summary_files:
        files.append(os.path.join(summary_dir, i))
    return files


def display_prompt_content(output_path, file_name, instruction, prompt, timestamp=False, join_sep="\n"):
    print("display_prompt_content", output_path, file_name)
    file_path = os.path.join(output_path, "subtitles", file_name)
    if join_sep is None or join_sep == "None":
        join_sep = ""
    elif join_sep == "\\n":
        join_sep = "\n"
    elif join_sep == "\\t":
        join_sep = "\t"
    elif join_sep == "\" \"":
        join_sep = " "
    text = get_formatted_subtitles(file_path, timestamp, join_sep)
    messages = f"{prompt}\n\n{text}"
    return messages


def display_file_content(output_path, file_name):
    file_path = os.path.join(output_path, "subtitles", file_name)
    with open(os.path.join(file_path), "r") as f:
        content = f.read()
    frames_dir = os.path.join(output_path, "frames")
    image_files = sorted([f for f in os.listdir(frames_dir) if f.endswith((".jpg", ".png", ".jpeg"))],
                         key=lambda x: int(os.path.splitext(x)[0]))

    image_file = image_files[int(os.path.splitext(file_name)[0])]
    image_file = os.path.join(frames_dir, image_file)
    return content, image_file


def extract_relative_subtitle(video_path, output_extract_folder):
    # 计算frames的时间段
    time_segments = calculate_time_segments(video_path, output_extract_folder)
    output_subtitle_path = os.path.join(output_extract_folder, "subtitles")
    if not os.path.exists(output_subtitle_path):
        os.makedirs(output_subtitle_path)
    output_audio_file = os.path.join(output_extract_folder, "audio.mp3")
    output_subtitle_file = os.path.join(output_extract_folder, "subtitle.txt")
    if not os.path.exists(output_audio_file):
        print("extracted video audio use moviepy")
        extract_audio_to_file(video_path, output_audio_file)
    else:
        print("extracted video audio use existed audio file")
    extracted_audio_text = ""
    if not os.path.exists(output_subtitle_file):
        print("extracted video subtitle use asr model")
        # save audio subtitles
        audio_text_result = audio_to_text(output_audio_file)
        extracted_audio_text = format_audio_chunks(audio_text_result["chunks"])
        save_to_file(output_subtitle_file, extracted_audio_text)
    else:
        print("extracted video subtitle use existed subtitle file")
        extracted_audio_text = read_from_file(output_subtitle_file)
    extract_relative_subtitle_segment(extracted_audio_text, output_extract_folder, time_segments)


# 定义一个函数来将时间字符串转换为秒数
def time_to_seconds(time_str):
    hours, minutes, seconds = map(int, time_str.split(':'))
    return hours * 3600 + minutes * 60 + seconds


def split_text_into_chunks(input_data):
    # 解析输入数据
    output = []
    for line in input_data.strip().split('\n'):
        match = re.match(r'(\d{2}:\d{2}:\d{2})--(\d{2}:\d{2}:\d{2})\s+(.*)', line)
        if match:
            start_time, end_time, text = match.groups()
            start_seconds = time_to_seconds(start_time)
            end_seconds = time_to_seconds(end_time)
            output.append({
                "timestamp": (start_seconds, end_seconds),
                "text": text
            })
    return output


def force_empty_directory(directory_path):
    if os.path.exists(directory_path):
        shutil.rmtree(directory_path)  # 删除目录及其所有内容
    os.makedirs(directory_path)  # 创建新的目录


def extract_relative_subtitle_segment(extracted_audio_text, output_extract_folder, time_segments):
    output_subtitle_path = os.path.join(output_extract_folder, "subtitles")
    force_empty_directory(output_subtitle_path)
    audio_text_result = split_text_into_chunks(extracted_audio_text)
    for i in range(len(time_segments)):
        audio_text_file = os.path.join(output_subtitle_path, f"{i}.txt")
        print(
            f"extract video {format_time(time_segments[i][0])}--{format_time(time_segments[i][1])} subtitle to file {audio_text_file}")
        filtered_subtitles = filter_subtitles(time_segments[i], audio_text_result)
        with open(audio_text_file, "w") as f:
            f.writelines(format_audio_chunks(filtered_subtitles))


def filter_subtitles(time_segments, chunks):
    """
    Filter the subtitle chunks that fall within a given time range.

    :param time_segments: A tuple (start_time, end_time) defining the time range.
    :param chunks: A list of dictionaries, where each dictionary has a 'timestamp' key with a tuple (start_time, end_time).
    :return: A list of chunks that fall within the specified time range.
    """
    filtered_chunks = []
    start_segment, end_segment = time_segments

    for chunk in chunks:
        start_time, end_time = chunk['timestamp']
        # Check if the chunk's time overlaps with the time segment
        if start_time < end_segment and end_time > start_segment:
            filtered_chunks.append(chunk)

    return filtered_chunks


def extract_audio_to_file(video_path, output_file):
    video = VideoFileClip(video_path)
    audio = video.audio
    audio.write_audiofile(output_file)
    video.close()


def extract_audio_segment(video_path, start_time, end_time, output_audio_path):
    """
    Extracts an audio segment from a video file.

    :param video_path: Path to the video file.
    :param start_time: Start time in seconds.
    :param end_time: End time in seconds.
    :param output_audio_path: Path to save the extracted audio file.
    """
    # Load the video file
    video = VideoFileClip(video_path)
    # Extract the segment
    audio_segment = video.subclip(start_time, end_time).audio
    # Write the audio segment to a file
    audio_segment.write_audiofile(output_audio_path)


def clean_cuda_memory():
    import gc
    gc.collect()
    torch.cuda.empty_cache()


class ModelContext:
    def __init__(self, model_id, model_type):
        self.model_id = model_id
        self.model_type = model_type
        self.model = None
        self.processor = None
        self.pipeline = None
        self.tokenizer = None

    def __enter__(self):
        device = "cuda:0" if torch.cuda.is_available() else "cpu"
        torch_dtype = torch.float16 if torch.cuda.is_available() else torch.float32

        if self.model_type == "audio":
            self.model = AutoModelForSpeechSeq2Seq.from_pretrained(
                self.model_id, torch_dtype=torch_dtype, low_cpu_mem_usage=False, use_safetensors=True
            )
            self.model.to(device)
            self.processor = AutoProcessor.from_pretrained(self.model_id)
            self.pipeline = pipeline(
                "automatic-speech-recognition",
                model=self.model,
                tokenizer=self.processor.tokenizer,
                feature_extractor=self.processor.feature_extractor,
                max_new_tokens=128,
                chunk_length_s=30,
                batch_size=16,
                return_timestamps=True,
                torch_dtype=torch_dtype,
                device=device,
            )
        elif self.model_type == "llm":
            self.model = AutoModelForCausalLM.from_pretrained(
                self.model_id,
                torch_dtype=torch_dtype,
                device_map=device,
                # load_in_8bit=True,  # 使用8位量化（可选）
            )
            self.tokenizer = AutoTokenizer.from_pretrained(self.model_id)

        return self

    def __exit__(self, exc_type, exc_value, traceback):
        if self.model is not None:
            del self.model
        if self.processor is not None:
            del self.processor
        if self.pipeline is not None:
            del self.pipeline
        if self.tokenizer is not None:
            del self.tokenizer
        clean_cuda_memory()


def format_audio_chunks(chunks, timestamp=True):
    formatted_text = ""
    for chunk in chunks:
        if timestamp:
            text = f"{format_time(chunk['timestamp'][0])}--{format_time(chunk['timestamp'][1])} {chunk['text']}\n"
        else:
            text = f"{chunk['text']}\n"
        formatted_text = formatted_text + text
    return formatted_text


def audio_to_text(audio_file):
    global global_llm_model_context
    if global_llm_model_context:
        global_llm_model_context.__exit__(None, None, None)
        global_llm_model_context = None
    with ModelContext("openai/whisper-large-v3", "audio") as ctx:
        result = ctx.pipeline(audio_file)
    return result


global_llm_model_context, llm_model_context = None, None


def generate_text(instruction, prompt, max_new_tokens=512, temperature=0.2, top_p=0.9, top_k=20):
    global global_llm_model_context, llm_model_context
    if global_llm_model_context is None:
        global_llm_model_context = ModelContext("Qwen/Qwen2-7B-Instruct", "llm")
        llm_model_context = global_llm_model_context.__enter__()
    messages = [
        {"role": "system", "content": instruction},
        {"role": "user", "content": prompt}
    ]
    text = llm_model_context.tokenizer.apply_chat_template(
        messages,
        tokenize=False,
        add_generation_prompt=True
    )
    try:
        with torch.no_grad():
            model_inputs = llm_model_context.tokenizer([text], return_tensors="pt").to(device)
            generated_ids = llm_model_context.model.generate(
                model_inputs.input_ids,
                max_new_tokens=max_new_tokens,
                temperature=temperature,
                top_p=top_p,
                top_k=top_k,
            )
            generated_ids = [
                output_ids[len(input_ids):] for input_ids, output_ids in
                zip(model_inputs.input_ids, generated_ids)
            ]
            response = llm_model_context.tokenizer.batch_decode(generated_ids, skip_special_tokens=True)[0]
    except Exception as e:
        print(e)
        response = ""
    return response


def save_to_file(file_name, data):
    with open(file_name, "w") as f:
        f.write(data)


def read_from_file(file_name):
    with open(file_name, "r") as f:
        return f.read()


def get_all_subtitles(output_extract_folder, timestamp=False, join_sep="\n"):
    subtitles = []
    subtitles_dir = os.path.join(output_extract_folder, "subtitles")
    subtitle_files = sorted([f for f in os.listdir(subtitles_dir) if f.endswith(".txt")],
                            key=lambda x: int(os.path.splitext(x)[0]))

    for index in range(len(subtitle_files)):
        subtitle_file = os.path.join(subtitles_dir, subtitle_files[index])
        subtitle = get_formatted_subtitles(subtitle_file, timestamp, join_sep)
        subtitles.append(subtitle)
    return subtitles


def get_formatted_subtitles(subtitle_file, timestamp=False, join_sep="\n"):
    subtitle = ""
    with open(subtitle_file, "r") as f:
        subtitle_lines = f.readlines()
        if not timestamp:
            for line in subtitle_lines:
                match = re.match(r'(\d{2}:\d{2}:\d{2})--(\d{2}:\d{2}:\d{2})\s+(.*)', line)
                if match:
                    start_time, end_time, text = match.groups()
                    subtitle = f"{subtitle}{join_sep}{text}"
        else:
            subtitle = join_sep.join(subtitle_lines)
    return subtitle


def summarize_relative_subtitle(instruction, prompt, max_new_tokens=512, temperature=0.2, top_p=0.9, top_k=20):
    response = generate_text(instruction, prompt, max_new_tokens, temperature,
                             top_p, top_k)
    return response


def summarize_all_relative_subtitle(output_extract_folder, instruction, prompt, timestamp=False, join_sep="\n",
                                    max_new_tokens=512, temperature=0.2,
                                    top_p=0.9, top_k=20):
    summary_path = os.path.join(output_extract_folder, "summary")
    force_empty_directory(summary_path)
    # subtitles 从文件中读取
    subtitles = get_all_subtitles(output_extract_folder, timestamp, join_sep)
    for i in range(len(subtitles)):
        subtitle = subtitles[i]
        summary_file = os.path.join(summary_path, f"{i}.md")
        print(f"extract video {i} subtitle summary to file {summary_file}")
        response = llm_subtitle_summary(subtitle, summary_file, instruction, prompt, max_new_tokens,
                                        temperature, top_p, top_k)
        save_to_file(summary_file, response)
    return list_all_summary_files(output_extract_folder)


def summarize_video_summary(output_extract_folder, instruction, prompt, max_new_tokens=512, temperature=0.2,
                            top_p=0.9, top_k=20):
    summary_folder = os.path.join(output_extract_folder, "summary")
    summary_files = sorted(
        [f for f in os.listdir(summary_folder) if f.endswith(".md") and f != "summary.md" and f != "mindmap.md"],
        key=lambda x: int(os.path.splitext(x)[0]))
    all_response = "```\n"
    for i in range(len(summary_files)):
        response = read_from_file(os.path.join(summary_folder, summary_files[i]))
        all_response = f"{all_response}第{i + 1}页\n\n{response}\n"
    all_response = all_response + "```"
    try:
        return llm_video_summary(all_response, summary_folder, instruction, prompt, max_new_tokens, temperature, top_p,
                                 top_k)
    except Exception as e:
        error_details = traceback.format_exc()
        print(error_details)
        return error_details


def summarize_video_summary_preview(output_extract_folder, instruction, prompt):
    summary_folder = os.path.join(output_extract_folder, "summary")
    summary_files = sorted(
        [f for f in os.listdir(summary_folder) if f.endswith(".md") and f != "summary.md" and f != "mindmap.md"],
        key=lambda x: int(os.path.splitext(x)[0]))
    all_response = "```\n"
    for i in range(len(summary_files)):
        response = read_from_file(os.path.join(summary_folder, summary_files[i]))
        all_response = f"{all_response}第{i + 1}页\n\n{response}\n"
    all_response = all_response + "```"
    message = f"Instruction:\n{instruction}\n\nUser:\n{prompt}\n\n{all_response}"
    return message


def summarize_video_mindmap(output_extract_folder, instruction, prompt, max_new_tokens=512, temperature=0.2,
                            top_p=0.9, top_k=20):
    summary_folder = os.path.join(output_extract_folder, "summary")
    summary_files = sorted(
        [f for f in os.listdir(summary_folder) if f.endswith(".md") and f != "summary.md" and f != "mindmap.md"],
        key=lambda x: int(os.path.splitext(x)[0]))
    all_response = "```\n"
    for i in range(len(summary_files)):
        response = read_from_file(os.path.join(summary_folder, summary_files[i]))
        all_response = f"{all_response}第{i + 1}页\n\n{response}\n"
    all_response = all_response + "```"
    try:
        return llm_video_mindmap(all_response, summary_folder, instruction, prompt, max_new_tokens, temperature,
                                 top_p, top_k)
    except Exception as e:
        error_details = traceback.format_exc()
        print(error_details)
        return error_details


def llm_subtitle_summary(subtitle, summary_file, instruction, prompt, max_new_tokens=512, temperature=0.2,
                         top_p=0.9, top_k=20):
    prompt = f"{prompt} \n" + subtitle
    response = generate_text(instruction, prompt, max_new_tokens, temperature,
                             top_p, top_k)
    save_to_file(summary_file, response)
    return response


def llm_video_summary(all_response, summary_path, instruction, prompt, max_new_tokens=512, temperature=0.2,
                      top_p=0.9, top_k=20):
    prompt = f"{prompt} \n\n" + all_response
    save_to_file(os.path.join(summary_path, f"summary_prompt.txt"), prompt)
    response = generate_text(instruction, prompt, max_new_tokens, temperature,
                             top_p, top_k)
    summary_file = os.path.join(summary_path, f"summary.md")
    print(f"extract video subtitle summary to file {summary_file}")
    save_to_file(summary_file, response)
    return response


def llm_video_mindmap(all_response, summary_path, instruction, prompt, max_new_tokens=512, temperature=0.2,
                      top_p=0.9, top_k=20):
    prompt = "以下是视频每一页ppt的总结，请你帮我尝试总结整个视频的主要内容，生成思维导图，请使用markdown格式进行输出该思维导图，请用中文回答 \n" + all_response
    save_to_file(os.path.join(summary_path, f"mindmap_prompt.txt"), prompt)
    response = generate_text(instruction, prompt, max_new_tokens, temperature,
                             top_p, top_k)
    mindmap_file = os.path.join(summary_path, f"mindmap.md")
    print(f"extract video subtitle summary to file {mindmap_file}")
    save_to_file(mindmap_file, response)
    return response


# 将Markdown内容转换为HTML
def markdown_to_html(markdown_content):
    html_content = markdown.markdown(markdown_content)
    return html_content


# 将html内容转为pdf
def html_to_pdf(html_content, output_path):
    import pdfkit
    pdfkit.from_string(html_content, output_path)


def parse_markdown(md_text):
    md = markdown_it.MarkdownIt()
    tokens = md.parse(md_text)
    return tokens


def create_mind_map(tokens):
    from graphviz import Digraph
    dot = Digraph(comment='Mind Map')
    stack = []
    for token in tokens:
        if token.type == 'heading_open':
            level = int(token.tag[1])
            while stack and stack[-1][0] >= level:
                stack.pop()
            stack.append((level, token.map[0]))
        elif token.type == 'inline' and stack:
            label = token.content.strip()
            if label:
                current_level, start_line = stack[-1]
                node_id = f'{start_line}_{current_level}'
                dot.node(node_id, label)
                if len(stack) > 1:
                    parent_level, parent_start_line = stack[-2]
                    parent_node_id = f'{parent_start_line}_{parent_level}'
                    dot.edge(parent_node_id, node_id)


def generate_mindmap_image(md_text):
    tokens = parse_markdown(md_text)
    mind_map = create_mind_map(tokens)
    mind_map.render('mind_map', format='png', view=True)


def delete_small_files(frames_dir, subtitles_dir, file_size_threshold=1024):
    # 获取subtitle目录中的所有文件
    subtitle_files = sorted([f for f in os.listdir(subtitles_dir) if f.endswith(".txt")],
                            key=lambda x: int(os.path.splitext(x)[0]))
    image_files = sorted([f for f in os.listdir(frames_dir) if f.endswith((".jpg", ".png", ".jpeg"))],
                         key=lambda x: int(os.path.splitext(x)[0]))

    for index in range(len(subtitle_files)):
        subtitle_file = subtitle_files[index]
        subtitle_path = os.path.join(subtitles_dir, subtitle_file)
        # 检查文件大小
        if os.path.getsize(subtitle_path) < 1024:  # 1KB = 1024 bytes
            # 获取对应的frame文件名
            frame_file = os.path.join(frames_dir, image_files[index])

            if os.path.exists(frame_file):
                # 删除对应的frame文件
                os.remove(frame_file)
                print(f"Deleted frame file: {frame_file}")
                # 删除subtitle文件
                os.remove(subtitle_path)
                print(f"Deleted subtitle file: {subtitle_path}")
            else:
                print(f"Frame file not found: {frame_file}")


def merge_small_subtitles(video_path, output_extract_folder, file_size_threshold=1024):
    frames_dir = os.path.join(output_extract_folder, "frames")
    subtitles_dir = os.path.join(output_extract_folder, "subtitles")
    delete_small_files(frames_dir, subtitles_dir, file_size_threshold)
    extract_relative_subtitle(video_path, output_extract_folder)


def images_to_ppt(output_extract_folder):
    from pptx import Presentation
    from pptx.util import Inches
    output_ppt = os.path.join(output_extract_folder, "output.pptx")
    presentation = Presentation()
    frame_path = os.path.join(output_extract_folder, "frames")
    image_files = sorted([f for f in os.listdir(frame_path) if f.endswith((".jpg", ".png", ".jpeg"))],
                         key=lambda x: int(os.path.splitext(x)[0]))

    for i, filename in enumerate(image_files):
        # Add image slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        img_path = os.path.join(frame_path, filename)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=presentation.slide_width)

        # Add a slide with time segment
        text_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        txBox = text_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tf = txBox.text_frame
        # p1 = tf.add_paragraph()
        # p1.text = f"Video Time: {format_time(time_segments[i][0])} to {format_time(time_segments[i][1])}"
        p2 = tf.add_paragraph()
        with open(os.path.join(output_extract_folder, "summary", f"{i}.md"), "r") as f:
            p2.text = f.read()
    presentation.save(output_ppt)
    print(f"PPT saved at: {output_ppt}")


def merge_frames_and_summary(extracted_folder, output_md):
    frames_folder = os.path.join(extracted_folder, "frames")
    summary_folder = os.path.join(extracted_folder, "summary")

    # 获取所有图片文件并按名称排序
    image_files = sorted([f for f in os.listdir(frames_folder) if f.endswith((".jpg", ".jpeg", ".png"))],
                         key=lambda x: int(os.path.splitext(x)[0]))

    # 获取所有Markdown文件并按名称排序
    summary_files = sorted(
        [f for f in os.listdir(summary_folder) if f.endswith(".md") and f != "summary.md" and f != "mindmap.md"],
        key=lambda x: int(os.path.splitext(x)[0]))

    # 读取mindmap.md和summary.md内容
    with open(os.path.join(summary_folder, "mindmap.md"), "r") as f:
        mindmap_content = f.read()

    with open(os.path.join(summary_folder, "summary.md"), "r") as f:
        summary_content = f.read()

    # 创建最终的Markdown文件
    with open(output_md, "w") as output_file:
        # 写入mindmap.md内容
        output_file.write(mindmap_content)
        output_file.write("\n\n")

        # 遍历图片和Markdown文件，按顺序写入
        for i, image_file in enumerate(image_files):
            # 写入图片
            output_file.write(f"![{image_file}](frames/{image_file})\n\n")

            # 写入对应的Markdown内容
            with open(os.path.join(summary_folder, summary_files[i]), "r") as f:
                summary_text = f.read()
                output_file.write(summary_text)
                output_file.write("\n\n")

        # 写入summary.md内容
        output_file.write(summary_content)

    print(f"Markdown file saved at: {output_md}")


def zip_files(output_extract_folder, output_file, file_list):
    """
    将指定的文件和目录列表压缩成一个ZIP文件。

    :param output_extract_folder: 包含文件和目录路径
    :param output_file: 输出ZIP文件的路径
    """
    with zipfile.ZipFile(output_file, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for item in file_list:
            if os.path.isfile(item):
                # 如果是文件，直接添加到ZIP文件
                zipf.write(item, os.path.basename(item))
            elif os.path.isdir(item):
                # 如果是目录，递归添加目录中的所有文件
                for root, dirs, files in os.walk(item):
                    for file in files:
                        file_path = os.path.join(root, file)
                        arcname = os.path.relpath(file_path, start=os.path.dirname(item))
                        zipf.write(file_path, arcname)
            else:
                print(f"警告：{item} 既不是文件也不是目录，跳过。")
    print(f"save output_extract_folder to {output_file} zip file done")


def save_output_to_file(output_extract_folder, subtitle, seg_subtitle, seg_summary):
    output_md = os.path.join(output_extract_folder, "movie.md")
    merge_frames_and_summary(output_extract_folder, output_md)
    images_to_ppt(output_extract_folder)
    file_list = [
        os.path.join(output_extract_folder, "frames"),
        os.path.join(output_extract_folder, "summary"),
        os.path.join(output_extract_folder, "output.pptx"),
        os.path.join(output_extract_folder, "movie.md"),
    ]
    if subtitle:
        file_list.append(os.path.join(output_extract_folder, "subtitle.txt"))
    if seg_subtitle:
        file_list.append(os.path.join(output_extract_folder, "subtitles"))
    if seg_summary:
        file_list.append(os.path.join(output_extract_folder, "summary"))
    output_zip_file = os.path.join(output_extract_folder, "output.zip")
    zip_files(output_extract_folder, output_zip_file, file_list)
    return output_zip_file


def main(video_file, output_extract_folder, interval_sec, merge_seg=True):
    output_frame_path = os.path.join(output_extract_folder, "frames")
    # Assume extract_similar_frames is defined elsewhere
    if not os.path.exists(output_frame_path) or not os.listdir(output_frame_path):
        extract_similar_frames(video_file, output_extract_folder, interval_sec)

    # Ensure the frames directory exists and contains images
    if not os.path.exists(output_frame_path) or not os.listdir(output_frame_path):
        print("No frames found. Please check your video processing.")
        return

    # 调用asr模型输出字幕，并将字幕按照frame分片
    extract_relative_subtitle(video_file, output_extract_folder)

    # 将subtitle中的小文件删除，同时删除对应的frames，然后重新生成subtitle来实现merge功能
    if merge_seg:
        merge_small_subtitles(video_file, output_extract_folder)

    # 调用llm对每个分段进行总结
    summarize_relative_subtitle(output_extract_folder)

    # 将以上内容输出到一个ppt文件中
    images_to_ppt(output_extract_folder)

    print("Process completed!")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract frames from video, create PPT with time info slides")
    parser.add_argument("--video_path", type=str, required=True, help="Path to the input video")
    parser.add_argument("--output_path", type=str, required=True, help="Path for the output PPT file")
    parser.add_argument("--interval_sec", type=float, default=1.5, help="Interval in seconds to capture frames")
    parser.add_argument("--merge_seg", type=bool, default=True, help="Merge small subtitle")

    args = parser.parse_args()
    OUTPUT_EXTRACT_FOLDER = args.output_path
    if not os.path.exists(OUTPUT_EXTRACT_FOLDER):
        os.makedirs(OUTPUT_EXTRACT_FOLDER)
    VIDEO_FILE = args.video_path
    if not os.path.exists(args.video_path):
        print("Error: Video file does not exist")
        import sys

        sys.exit(1)

    device = "cuda:0" if torch.cuda.is_available() else "cpu"

    try:
        main(args.interval_sec, args.merge_seg)
    except Exception as e:
        print(f"Error during processing: {str(e)}")
